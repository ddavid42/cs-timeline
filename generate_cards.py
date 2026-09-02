import csv
import copy
import os
import argparse

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Pt
from PIL import Image, ImageDraw, ImageFont


# ============================================================
# DEFAULT CONFIGURATION
# ============================================================

DEFAULT_TEMPLATE = "template.pptx"
DEFAULT_CSV = "text_card_fr.csv"
DEFAULT_IMAGE_DIR = "images"
DEFAULT_OUTPUT = "card_fr.pptx"


# ============================================================
# POWERPOINT UTILITIES
# ============================================================

def find_shape(slide, name):
    """
    Find a shape by its PowerPoint name.

    Example:
        find_shape(slide, "IMG_FOND")
    """

    for shape in slide.shapes:
        if shape.name == name:
            return shape

    available = [shape.name for shape in slide.shapes]

    raise ValueError(
        f"\nObject '{name}' not found on the slide.\n"
        f"Available objects: {available}\n"
    )


def replace_image(shape, image_path):
    """
    Replace the image in an existing shape.

    The shape's position, size, crop, and other properties are preserved.
    """

    if not os.path.isfile(image_path):
        raise FileNotFoundError(
            f"Image not found: {image_path}"
        )

    # Get the existing image relationship.
    blip = shape._element.blipFill.blip

    # Add the new image to the PowerPoint package.
    image_part, rId = shape.part.get_or_add_image_part(
        image_path
    )

    # Replace the reference.
    blip.set(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed",
        rId
    )


def set_run_text_with_line_breaks(run, text):
    """
    Set run text, converting "\n" to PowerPoint line breaks (<a:br/>)
    within the same paragraph.

    The run formatting (font, size, bold, and so on) is copied to each
    inserted line.
    """

    lines = text.split("\n")

    run.text = lines[0]

    anchor = run._r

    for line in lines[1:]:

        break_element = anchor.makeelement(qn("a:br"), {})

        run_properties = anchor.find(qn("a:rPr"))

        if run_properties is not None:
            break_element.append(copy.deepcopy(run_properties))

        anchor.addnext(break_element)

        new_run_element = copy.deepcopy(anchor)
        new_run_element.find(qn("a:t")).text = line

        break_element.addnext(new_run_element)

        anchor = new_run_element


def find_placeholder_run(shape, placeholder):
    """
    Return the (run, paragraph) containing the requested placeholder text
    exactly (for example, "TITLE" or "DATE").
    """

    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text == placeholder:
                return run, paragraph

    raise ValueError(
        f"Placeholder text '{placeholder}' was not found "
        f"in shape '{shape.name}'."
    )


def replace_placeholder_text(shape, placeholder, new_text):
    """
    Replace a run in a shape whose text exactly matches a placeholder
    (for example, "TITLE" or "DATE").

    The run formatting (font, size, bold, and so on) is preserved.
    A "\n" in new_text becomes a real PowerPoint line break.
    """

    # "\x0b" is PowerPoint's own line break, distinct from "\n". It is not
    # valid XML; lxml would escape it to the visible string "_x000B_".
    new_text = new_text.replace("\x0b", " ")

    run, _ = find_placeholder_run(shape, placeholder)
    set_run_text_with_line_breaks(run, new_text)


def remove_placeholder_paragraph(shape, placeholder):
    """
    Remove the entire paragraph containing a placeholder. This is used to
    remove an empty subtitle and vertically center the remaining text.
    """

    _, paragraph = find_placeholder_run(shape, placeholder)
    paragraph._p.getparent().remove(paragraph._p)


# ============================================================
# AUTOMATIC FONT-SIZE ADJUSTMENT
# ============================================================
#
# LibreOffice, unlike PowerPoint, does not apply automatic text shrinking on
# overflow. The following helpers estimate the wrapped text and calculate the
# largest font size that fits in the shape without going below the minimum.

CONTENT_MIN_FONT_SIZE_PT = 8
CONTENT_FONT_SIZE_STEP_PT = 0.5
LINE_SPACING_FACTOR = 1.2

_MEASURING_DRAW = ImageDraw.Draw(Image.new("RGB", (1, 1)))
_MEASURING_FONT_CACHE = {}


def _get_measuring_font(size_pt):
    """Return an approximate font used only to estimate text widths."""

    size_px = max(1, round(size_pt))

    font = _MEASURING_FONT_CACHE.get(size_px)

    if font is None:
        font = ImageFont.load_default(size=size_px)
        _MEASURING_FONT_CACHE[size_px] = font

    return font


def _wrap_line(text, font, max_width_pt):
    """Wrap one line of text to the available width."""

    if not text:
        return [""]

    words = text.split(" ")
    lines = []
    current = ""

    for word in words:
        candidate = f"{current} {word}".strip()

        if current and _MEASURING_DRAW.textlength(candidate, font=font) > max_width_pt:
            lines.append(current)
            current = word
        else:
            current = candidate

    lines.append(current)

    return lines


def _count_wrapped_lines(text, font_size_pt, max_width_pt):
    font = _get_measuring_font(font_size_pt)

    return sum(
        len(_wrap_line(segment, font, max_width_pt))
        for segment in text.split("\n")
    )


def _content_text_fits(shape, title, content_text, title_size_pt, content_size_pt):
    text_frame = shape.text_frame

    usable_width_pt = (
        shape.width - text_frame.margin_left - text_frame.margin_right
    ) / 12700

    usable_height_pt = (
        shape.height - text_frame.margin_top - text_frame.margin_bottom
    ) / 12700

    title_lines = _count_wrapped_lines(title, title_size_pt, usable_width_pt)
    content_lines = _count_wrapped_lines(content_text, content_size_pt, usable_width_pt)

    needed_height_pt = LINE_SPACING_FACTOR * (
        title_lines * title_size_pt
        + content_size_pt  # ligne vide de séparation entre titre et texte
        + content_lines * content_size_pt
    )

    return needed_height_pt <= usable_height_pt


def fit_content_font_size(shape, title, content_text, base_size_pt, title_size_pt):
    """
    Return the largest font size, in points, between CONTENT_MIN_FONT_SIZE_PT
    and base_size_pt that lets the text fit inside the shape.
    """

    size = base_size_pt

    while size > CONTENT_MIN_FONT_SIZE_PT:

        if _content_text_fits(shape, title, content_text, title_size_pt, size):
            return size

        size = round(size - CONTENT_FONT_SIZE_STEP_PT, 1)

    return CONTENT_MIN_FONT_SIZE_PT


R_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"


def duplicate_slide(prs, source_slide):
    """
    Duplicate an existing slide.

    All template shapes and the relationships they reference (images, etc.)
    are copied.
    """

    new_slide = prs.slides.add_slide(
        source_slide.slide_layout
    )

    # Remove placeholders automatically added from the layout.
    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)

    # Copy relationships and map old rIds to new ones. An rId cannot be reused
    # unchanged from one slide part to another.
    rid_map = {}

    for rId, rel in source_slide.part.rels.items():

        if "notesSlide" in rel.reltype or "slideLayout" in rel.reltype:
            continue

        target = rel.target_ref if rel.is_external else rel.target_part

        new_rid = new_slide.part.relate_to(
            target,
            rel.reltype,
            is_external=rel.is_external
        )

        rid_map[rId] = new_rid

    # Copy shapes while remapping their relationship references.
    for shape in source_slide.shapes:

        new_shape = copy.deepcopy(
            shape.element
        )

        for attr in ("embed", "link", "id"):
            for element in new_shape.iter():
                old_rid = element.get(R_NS + attr)

                if old_rid in rid_map:
                    element.set(R_NS + attr, rid_map[old_rid])

        new_slide.shapes._spTree.insert_element_before(
            new_shape,
            "p:extLst"
        )

    return new_slide



def delete_slide(prs, slide):
    """
    Delete a slide from the presentation.
    """

    slide_id = slide.slide_id

    for slide_id_element in prs.slides._sldIdLst:

        if int(slide_id_element.id) == slide_id:

            rId = slide_id_element.rId

            prs.part.drop_rel(rId)

            prs.slides._sldIdLst.remove(
                slide_id_element
            )

            break


# ============================================================
# CARD CONFIGURATION
# ============================================================

def configure_front(slide, image_path, title, subtitle):
    """
    Configure a card front.
    """

    # Background image
    background = find_shape(
        slide,
        "IMG_FOND"
    )

    replace_image(
        background,
        image_path
    )

    # Title / subtitle
    title_shape = find_shape(
        slide,
        "TITRE"
    )

    replace_placeholder_text(title_shape, "TITLE", title)

    if subtitle:
        replace_placeholder_text(title_shape, "SUB-TITLE", subtitle)
    else:
        # Remove the subtitle paragraph to vertically center the title.
        remove_placeholder_paragraph(title_shape, "SUB-TITLE")


def configure_back(
    slide,
    image_path,
    title,
    date,
    text
):
    """
    Configure a card back.
    """

    # --------------------------------------------------------
    # Background image
    # --------------------------------------------------------

    background = find_shape(
        slide,
        "IMG_FOND"
    )

    replace_image(
        background,
        image_path
    )

    # --------------------------------------------------------
    # Content (title + explanatory text)
    # --------------------------------------------------------

    content = find_shape(
        slide,
        "CONTENU"
    )

    # Get the original font sizes before replacement, so the final text can be
    # checked for fit within the content box.
    title_run, _ = find_placeholder_run(content, "TITLE")
    content_run, content_paragraph = find_placeholder_run(content, "CONTENT")

    base_title_size_pt = title_run.font.size.pt if title_run.font.size else 14
    base_content_size_pt = content_run.font.size.pt if content_run.font.size else 10

    replace_placeholder_text(content, "TITLE", title)
    replace_placeholder_text(content, "CONTENT", text)

    # Reduce text font size automatically when it overflows the allocated box,
    # down to CONTENT_MIN_FONT_SIZE_PT.
    font_size_pt = fit_content_font_size(
        content,
        title,
        text,
        base_content_size_pt,
        base_title_size_pt
    )

    for run in content_paragraph.runs:
        run.font.size = Pt(font_size_pt)

    # --------------------------------------------------------
    # Date
    # --------------------------------------------------------

    date_shape = find_shape(
        slide,
        "DATE"
    )

    replace_placeholder_text(date_shape, "DATE", date)


# ============================================================
# TEMPLATE VALIDATION
# ============================================================

def validate_template(prs):
    """
    Check that the template contains the required objects.
    """

    if len(prs.slides) < 2:

        raise ValueError(
            "The PowerPoint template must contain at least two slides: "
            "front and back."
        )

    front = prs.slides[0]
    back = prs.slides[1]

    required_front = [
        "IMG_FOND",
        "IMG_OVERLAY",
        "TITRE"
    ]

    required_back = [
        "IMG_FOND",
        "IMG_OVERLAY",
        "CONTENU",
        "DATE"
    ]

    front_names = [
        shape.name
        for shape in front.shapes
    ]

    back_names = [
        shape.name
        for shape in back.shapes
    ]

    print("Validating template...")
    print()

    for name in required_front:

        if name not in front_names:

            raise ValueError(
                f"Object '{name}' is missing from the FRONT slide.\n"
                f"Available objects: {front_names}"
            )

    for name in required_back:

        if name not in back_names:

            raise ValueError(
                f"Object '{name}' is missing from the BACK slide.\n"
                f"Available objects: {back_names}"
            )

    print("  ✓ Front OK")
    print("  ✓ Back OK")
    print()


# ============================================================
# CSV READING
# ============================================================

def read_csv(csv_path):

    if not os.path.isfile(csv_path):

        raise FileNotFoundError(
            f"CSV file not found: {csv_path}"
        )

    cards = []

    with open(
        csv_path,
        "r",
        encoding="utf-8-sig",
        newline=""
    ) as file:

        reader = csv.DictReader(file)

        required = {
            "date",
            "title",
            "subtitle",
            "image",
            "text"
        }

        columns = set(
            reader.fieldnames or []
        )

        missing = required - columns

        if missing:

            raise ValueError(
                "Missing CSV columns: "
                + ", ".join(sorted(missing))
            )

        for line_number, row in enumerate(
            reader,
            start=2
        ):

            image = row["image"].strip()
            # The literal "\n" sequence in the CSV represents a manual line
            # break, which is easier to enter than an actual line break.
            title = row["title"].strip().replace("\\n", "\n")
            # The subtitle is optional and may be empty.
            subtitle = row["subtitle"].strip().replace("\\n", "\n")
            date = row["date"].strip()
            text = row["text"].strip().replace("\\n", "\n")

            if not image:
                raise ValueError(
                    f"Line {line_number}: "
                    "empty image value."
                )

            if not title:
                raise ValueError(
                    f"Line {line_number}: "
                    "empty title value."
                )

            if not date:
                raise ValueError(
                    f"Line {line_number}: "
                    "empty date value."
                )

            if not text:
                raise ValueError(
                    f"Line {line_number}: "
                    "empty text value."
                )

            cards.append({
                "image": image,
                "title": title,
                "subtitle": subtitle,
                "date": date,
                "text": text
            })

    return cards


# ============================================================
# GENERATION
# ============================================================

def generate_cards(
    template_path,
    csv_path,
    image_dir,
    output_path
):

    print("=" * 70)
    print("GENERATING CARDS")
    print("=" * 70)
    print()

    # --------------------------------------------------------
    # Checks
    # --------------------------------------------------------

    if not os.path.isfile(template_path):

        raise FileNotFoundError(
                f"PowerPoint template not found: "
            f"{template_path}"
        )

    cards = read_csv(csv_path)

    if not cards:

        raise ValueError(
            "The CSV does not contain any cards."
        )

    # --------------------------------------------------------
    # Load the template
    # --------------------------------------------------------

    prs = Presentation(
        template_path
    )

    validate_template(prs)

    # The first two slides are the templates.
    template_front = prs.slides[0]
    template_back = prs.slides[1]

    generated = []

    # --------------------------------------------------------
    # Generate cards
    # --------------------------------------------------------

    for number, card in enumerate(
        cards,
        start=1
    ):

        print(
            f"[{number}/{len(cards)}] "
            f"{card['title']}"
        )

        image_path = os.path.join(
            image_dir,
            card["image"]
        )

        if not os.path.isfile(image_path):

            raise FileNotFoundError(
                f"\nImage not found for "
                f"'{card['title']}':\n"
                f"{image_path}"
            )

        # ====================================================
        # FRONT
        # ====================================================

        front = duplicate_slide(
            prs,
            template_front
        )

        configure_front(
            front,
            image_path,
            card["title"],
            card["subtitle"]
        )

        generated.append(front)

        # ====================================================
        # BACK
        # ====================================================

        back = duplicate_slide(
            prs,
            template_back
        )

        configure_back(
            back,
            image_path,
            card["title"],
            card["date"],
            card["text"]
        )

        generated.append(back)

    # --------------------------------------------------------
    # Remove the two template slides.
    # --------------------------------------------------------

    delete_slide(
        prs,
        template_front
    )

    delete_slide(
        prs,
        template_back
    )

    # --------------------------------------------------------
    # Save
    # --------------------------------------------------------

    prs.save(
        output_path
    )

    print()
    print("=" * 70)
    print("GENERATION COMPLETE")
    print("=" * 70)
    print()
    print(f"Cards generated: {len(cards)}")
    print(f"Slides generated: {len(cards) * 2}")
    print()
    print(f"File: {output_path}")
    print()


# ============================================================
# MAIN PROGRAM
# ============================================================

def main():

    parser = argparse.ArgumentParser(
        description=(
            "Generate PowerPoint cards automatically from a template "
            "and a CSV file."
        )
    )

    parser.add_argument(
        "--template",
        default=DEFAULT_TEMPLATE,
        help="PowerPoint template file"
    )

    parser.add_argument(
        "--csv",
        default=DEFAULT_CSV,
        help="CSV file containing card data"
    )

    parser.add_argument(
        "--images",
        default=DEFAULT_IMAGE_DIR,
        help="Folder containing images"
    )

    parser.add_argument(
        "--output",
        default=DEFAULT_OUTPUT,
        help="Output PowerPoint file"
    )

    args = parser.parse_args()

    generate_cards(
        template_path=args.template,
        csv_path=args.csv,
        image_dir=args.images,
        output_path=args.output
    )


if __name__ == "__main__":

    main()